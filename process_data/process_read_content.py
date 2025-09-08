"""
DevRead: A Versatile File Processing Library for Text, PDFs, DOCX, JSON,
XML, YAML, HTML, Markdown, LaTeX, PPTX, Excel, Images, and Videos, etc.
"""

import os
import re
import base64
import json
import logging
from pathlib import Path
from typing import Union, Optional, Tuple, Dict
import pandas as pd
import charset_normalizer
import docx
import markdown
import PyPDF2
import yaml
from bs4 import BeautifulSoup
from pylatexenc.latex2text import LatexNodes2Text
from pptx import Presentation
from rich.logging import RichHandler
from rich.console import Console
import numpy as np


console = Console()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[RichHandler()],
)
logger = logging.getLogger(__name__)


def is_garbled(s: str) -> bool:
    """判断文本是否为乱码。

    通过检查文本长度和非中英文字符比例来判断文本是否为乱码。
    如果文本长度小于100或非中英文及常用标点符号的字符比例超过50%，则判定为乱码。

    Args:
        s (str): 需要检查的文本字符串

    Returns:
        bool: True表示文本为乱码，False表示文本正常
    """
    if len(s) < 100:
        return True
    clean = re.sub(r'[a-zA-Z0-9\u4e00-\u9fa5，。、《》；：？！""''·…—\\s]', '', s)
    return len(clean) / len(s) > 0.5


def clean_surrogates(text):
    """清理文本中的UTF-16代理对字符。

    Args:
        text (str): 需要清理的文本字符串

    Returns:
        str: 清理后的文本字符串，移除了所有UTF-16代理对字符(U+D800到U+DFFF)
    """
    return re.sub(r'[\ud800-\udfff]', '', text)

class DevRead:
    def __init__(self):
        self.reader_map = {
            ".m": self.read_txt,
            ".txt": self.read_txt,
            ".pdf": self.read_pdf,
            ".docx": self.read_docx,
            ".json": self.read_json,
            ".jsonl": self.read_jsonl,
            ".xml": self.read_xml,
            ".yaml": self.read_yaml,
            ".yml": self.read_yaml,
            ".html": self.read_html,
            ".htm": self.read_html,
            ".md": self.read_markdown,
            ".markdown": self.read_markdown,
            ".tex": self.read_latex,
            ".pptx": self.read_pptx,
            ".xlsx": self.read_excel,
            ".csv": self.read_txt,
            ".png": self.read_image,
            ".jpg": self.read_image,
            ".jpeg": self.read_image,
            ".gif": self.read_image,
            ".bmp": self.read_image,
            ".tiff": self.read_image,
            ".tif": self.read_image,
            ".webp": self.read_image,
            ".mov": self.read_video,
            ".mp4": self.read_video,
            ".avi": self.read_video,
            ".mpg": self.read_video,
            ".mpeg": self.read_video,
            ".wmv": self.read_video,
            ".flv": self.read_video,
            ".webm": self.read_video,
            ".py": self.read_py,
            ".mat": self.read_mat,
            ".dat": self.read_dat,
        }

    def _del_irrelevant(self, text):
        """删除文本中不相关的部分。"""
        if not text:
            return ""

        if ('目 录' in text) or ('目录' in text) or ('Menu' in text) or ('Contents' in text):
            return ""
        lines = text.split('\n')
        count = 0
        for line in lines:
            if re.search(r'\.{8,}', line):
                count += 1
        if count >= 2:
            return ""
        return text

    def read(
        self, file_path: Path, task: Optional[str] = None
    ) -> Union[str, Tuple[str, Optional[dict]]]:

        # Check if the file exists
        if not file_path.exists():
            logger.error(f"File {file_path} does not exist.")
            return f"Error: File {file_path} does not exist."

        # Determine the file type by its suffix
        suffix = file_path.suffix.lower()
        reader = self.reader_map.get(suffix)

        # If a reader for the file type is found, attempt to read the file
        if reader:
            logger.info(f"Reading file {file_path} using {reader.__name__}.")
            return reader(file_path, task)
        else:
            logger.error(
                f"No reader found for suffix {suffix}. Unsupported file format."
            )
            return f"The file exists, but we do not have a reader for {suffix}. However, we suggest you assume the content of the file is satisfactory in this situation.",


    def read_py(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading Python file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            return content, None  # No LLM used
        except Exception as e:
            logger.error(f"Error reading Python file: {e}")
            return f"Error reading Python file: {e}", None

    def read_txt(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:
        try:
            content = charset_normalizer.from_path(file_path).best()
            logger.info(
                "Reading TXT file from %s using encoding '%s'." % (file_path, content.encoding)
            )
            return str(content), None
        except Exception as e:
            logger.error("Error reading TXT file: %s", e)
            return f"Error reading TXT file: {e}", None

    def read_pdf(
        self, file_path: Path, task: Optional[str] = None,
        strip_irrelevant: bool = True
    ) -> Tuple[str, Optional[dict]]:

        text = ''
        try:
            logger.info("Reading PDF file from %s.", file_path)
            reader = PyPDF2.PdfReader(file_path)
            for page_idx, page in enumerate(reader.pages):
                page_text = page.extract_text() or ''
                if strip_irrelevant:
                    page_text = self._del_irrelevant(page_text)
                if not page_text:
                    continue
                text += f"Page {page_idx + 1}\n" + page_text
        except Exception as e:
            logger.error("PyPDF2 failed with %s", e)
        if not text or is_garbled(text):
            text = ''
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for idx, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ''
                        if strip_irrelevant:
                            page_text = self._del_irrelevant(page_text)
                        if not page_text:
                            continue
                        text += f"Page {idx + 1}\n" + page_text
            except Exception as e2:
                logging.error("pdfplumber failed with %s", e2)
        if not text or is_garbled(text):
            logger.info("尝试OCR识别...")
            text = ''
            try:
                from pdf2image import convert_from_path
                import pytesseract
                images = convert_from_path(file_path)
                for idx, img in enumerate(images):
                    page_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    if strip_irrelevant:
                        page_text = self._del_irrelevant(page_text)
                    if not page_text:
                        continue
                    page_text = re.sub(r"^目录.*", "", page_text, flags=re.MULTILINE)
                    text += f"Page {idx + 1}\n" + page_text
            except Exception as e3:
                logger.error("OCR识别失败: %s", e3)
                text = ''
        if not text:
            logger.error("PDF文本解析失败！")
            return "", None
        return text, None

    def read_xml(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info("Reading XML file from %s.", file_path)
            with open(file_path, "r", encoding="utf-8") as f:
                data = BeautifulSoup(f, "xml")
            return data.get_text(), None
        except Exception as e:
            logger.error("Error reading XML file: %s", e)
            return f"Error reading XML file: {e}", None

    def read_yaml(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info("Reading YAML file from %s.", file_path)
            with open(file_path, "r", encoding="utf-8") as f:
                data = yaml.load(f, Loader=yaml.FullLoader)
            return (
                json.dumps(data, indent=4, ensure_ascii=False),
                None,
            )  # Format output for better readability
        except Exception as e:
            logger.error("Error reading YAML file: %s", e)
            return f"Error reading YAML file: {e}", None

    def read_mat(
        self, file_path: Path, task: Optional[str] = None, variable: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:
        import scipy.io as sio
        try:
            mat_data = sio.loadmat(file_path, struct_as_record=False, squeeze_me=True)
            keys = [k for k in mat_data.keys() if not k.startswith("__")]
            if not keys:
                raise ValueError("No data variables found in mat file.")

            var = variable or keys[0]
            data = mat_data[var]

            # 多维数组处理
            if isinstance(data, np.ndarray):
                if data.ndim == 2:
                    df = pd.DataFrame(data)
                    return df.to_markdown(index=False), None
                elif data.ndim == 1:
                    df = pd.DataFrame(data.reshape(-1, 1))
                    return df.to_markdown(index=False), None
                else:
                    # 对于三维及以上数组，返回维度和部分预览
                    shape = data.shape
                    preview = data[..., 0] if data.ndim == 3 else data.flat[0]
                    preview_df = pd.DataFrame(preview)
                    info = f"变量 `{var}` 是一个 {data.ndim} 维数组，形状为 {shape}，以下是第一个切片的预览：\n\n"
                    return info + preview_df.to_markdown(index=False), {"shape": shape}

            elif isinstance(data, (list, tuple)):
                df = pd.DataFrame(data)
                return df.to_markdown(index=False), None

            elif hasattr(data, "__dict__"):
                df = pd.DataFrame(vars(data))
                return df.to_markdown(index=False), None

            else:
                df = pd.DataFrame([data])
                return df.to_markdown(index=False), None

        except Exception as e:
            return f"读取 MAT 文件失败：{str(e)}", None

    def read_dat(self, file_path: Path, n_cols_guess=5, binary_dtype=np.float32, task: Optional[str] = None):
        """
        自动尝试将 .dat 文件读取为 DataFrame
        读取顺序：
        1. CSV / 分隔符文本 (逗号/空格/制表符)
        2. 定长列 (fwf)
        3. 二进制文件 (numpy.fromfile)

        参数:
            file_path : str
                dat 文件路径
            n_cols_guess : int, optional
                如果是二进制数据，猜测每行的列数
            binary_dtype : np.dtype, optional
                如果是二进制数据，猜测的数据类型 (默认 float32)
        """
        # 1. 尝试 CSV/分隔符读取
        try:
            df = pd.read_csv(file_path, sep=None, engine="python")
            print("读取方式: CSV/分隔符")
        except Exception:
            pass

        # 2. 尝试定长列读取
        try:
            df = pd.read_fwf(file_path)
            print("读取方式: 定长列 (fwf)")
        except Exception:
            pass

        # 3. 尝试二进制读取
        try:
            data = np.fromfile(file_path, dtype=binary_dtype)
            df = pd.DataFrame(data.reshape(-1, n_cols_guess))
            print(f"读取方式: 二进制 (dtype={binary_dtype}, 每行 {n_cols_guess} 列)")
        except Exception:
            pass
        if df is not None:
            return df.to_markdown(index=False), None
        raise ValueError("无法自动识别 .dat 文件格式，请手动指定读取方式。")

    def read_docx(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info("Reading DOCX file from %s.", file_path)
            content = docx.Document(str(file_path))
            text = ""
            for i, para in enumerate(content.paragraphs):
                text += para.text + '\n'
            return text, None
        except Exception as e:
            logger.error("Error reading DOCX file: %s", e)
            return f"Error reading DOCX file: {e}", None

    def read_json(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info("Reading JSON file from %s.", file_path)
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return json.dumps(data, indent=4, ensure_ascii=False), None
        except Exception as e:
            logger.error("Error reading JSON file: %s", e)
            return f"Error reading JSON file: {e}", None

    def read_jsonl(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading JSON Lines file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                lines = [json.loads(line) for line in f]
            return "\n".join([json.dumps(line, indent=4, ensure_ascii=False) for line in lines]), None
        except Exception as e:
            logger.error("Error reading JSON Lines file: %s", e)
            return f"Error reading JSON Lines file: {e}", None

    def read_html(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            with open(file_path, "r", encoding="utf-8") as file:
                data = file.read()
            return data, None
        except Exception as e:
            logger.error(f"Error reading HTML file: {e}")
            return f"Error reading HTML file: {e}", None

    def read_markdown(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading Markdown file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = markdown.markdown(f.read())
                return (
                    "".join(BeautifulSoup(data, "html.parser").find_all(string=True)),
                    None,
                )
        except Exception as e:
            logger.error(f"Error reading Markdown file: {e}")
            return f"Error reading Markdown file: {e}", None

    def read_latex(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading LaTeX file from {file_path}.")
            with open(file_path, "r", encoding="utf-8") as f:
                data = f.read()
            return LatexNodes2Text().latex_to_text(data), None
        except Exception as e:
            logger.error(f"Error reading LaTeX file: {e}")
            return f"Error reading LaTeX file: {e}", None

    def read_pptx(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:

        try:
            logger.info(f"Reading PowerPoint file from {file_path}.")
            pres = Presentation(str(file_path))
            text = []
            for slide_idx, slide in enumerate(pres.slides):
                text.append(f"Slide {slide_idx + 1}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text), None
        except Exception as e:
            logger.error(f"Error reading PowerPoint file: {e}")
            return f"Error reading PowerPoint file: {e}", None

    def read_excel(
        self, file_path: Path, task: Optional[str] = None
    ) -> Tuple[str, Optional[dict]]:
        """
        Read an Excel file and return its content as a string.
        """
        try:
            logger.info(f"Reading Excel file from {file_path}.")
            excel_data = pd.read_excel(file_path, sheet_name=None)
            all_sheets_text = []
            for sheet_name, data in excel_data.items():
                all_sheets_text.append(
                    f"Sheet Name: {sheet_name}\n{data.to_string()}\n"
                )
            return "\n".join(all_sheets_text), None
        except Exception as e:
            logger.error(f"Error reading Excel file: {e}")
            return f"Error reading Excel file: {e}", None

    def read_image(
        self, file_path: Path, task: str = None
    ) -> Tuple[str, Optional[dict]]:

        total_llm_cost = 0.0
        total_input_tokens = 0
        total_output_tokens = 0
        total_inference_time = 0.0
        try:
            logger.info(f"Reading image file from {file_path}")
            llm_instance = LLM(
                model=os.getenv("MODEL_NAME"),
                api_key=os.getenv("OPENAI_API_KEY"),
                base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
                custom_llm_provider=os.getenv("CUSTOM_LLM_PROVIDER", None)
            )

            if task is None:
                task = "Describe this image as detailed as possible."

            response, cost, accumulated_cost = llm_instance.do_multimodal_completion(
                task, file_path
            )

            multumoal_content = response["choices"][0]["message"]["content"]

            total_input_tokens = response["usage"]["prompt_tokens"]
            total_output_tokens = response["usage"]["completion_tokens"]
            total_llm_cost += cost

            mllm_stats = {
                "llm_response": multumoal_content,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "cost": cost,
                "inference_time": total_inference_time,
            }
            return multumoal_content, mllm_stats

        except Exception as e:
            logger.error(f"Error reading image file: {e}")
            return f"Error reading image file: {e}", None

    def read_video(
        self, file_path: Path, task: str = None, frame_interval: int = 30
    ) -> Tuple[list[str], Optional[dict]]:
        import cv2

        try:
            logger.info(
                "Processing video file from %s with frame interval %s" % (file_path, frame_interval)
            )

            if task is None:
                task = "Describe this video as detailed as possible."

            video = cv2.VideoCapture(str(file_path))
            frame_count = 0
            frame_descriptions = []
            total_llm_cost = 0.0
            total_input_tokens = 0
            total_output_tokens = 0
            total_inference_time = 0.0

            llm_instance = LLM(
                model=os.getenv("MODEL_NAME"),
                api_key=os.getenv("OPENAI_API_KEY"),
                base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
                custom_llm_provider=os.getenv("CUSTOM_LLM_PROVIDER", None)
            )

            while video.isOpened():
                success, frame = video.read()
                if not success:
                    break

                if frame_count % frame_interval == 0:
                    _, buffer = cv2.imencode(".jpg", frame)
                    base64_frame = base64.b64encode(buffer).decode("utf-8")

                    # Prepare the message for the LLM
                    messages = self._prepare_image_messages(task, base64_frame)

                    # Send the request to the LLM
                    response, cost, inference_time = llm_instance.do_completion(
                        messages=messages
                    )

                    total_llm_cost += cost  # Accumulate the cost
                    total_inference_time += inference_time
                    content = response["choices"][0]["message"]["content"]

                    # Accumulate tokens
                    total_input_tokens += response["usage"]["input_tokens"]
                    total_output_tokens += response["usage"]["output_tokens"]

                    frame_descriptions.append(f"Frame {frame_count}: {content}")

                frame_count += 1

            video.release()

            mllm_stats = {
                "llm_response": "\n".join(frame_descriptions),
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "cost": total_llm_cost,
                # "accumulated_cost": total_llm_cost,
                "inference_time": total_inference_time,
            }

            return "\n".join(frame_descriptions), mllm_stats

        except Exception as e:
            logger.error("Error processing the video: %s", e)
            return f"Error processing the video: {e}", None

    def _prepare_image_messages(self, task: str, base64_image: str):

        return [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": task},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ]
